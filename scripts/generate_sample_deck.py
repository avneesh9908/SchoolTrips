"""
Builds the dummy "Parent's Orientation deck" referenced by the Documents tab.

Upload it to Drive, share it "Anyone with the link -> Viewer", and paste its link
into Documents!C2 so the site shows a real thumbnail instead of a placeholder.

Run:  python scripts/generate_sample_deck.py
"""

import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = pathlib.Path(__file__).resolve().parent.parent
OUT = ROOT / "sample-data" / "Grade 7 - Parent Orientation deck.pptx"

# Same palette as the website, so the deck and the site read as one thing.
BG = RGBColor(0xFC, 0xF8, 0xED)
INK = RGBColor(0x22, 0x30, 0x3F)
MUTED = RGBColor(0x76, 0x70, 0x66)
TEAL = RGBColor(0x26, 0xC0, 0xB0)   # Grade 7
RED = RGBColor(0xFF, 0x6B, 0x5B)
GREEN = RGBColor(0x4C, 0xAF, 0x6D)
AMBER = RGBColor(0xFF, 0xB1, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xEC, 0xE6, 0xD6)

W, H = Inches(13.333), Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False,
            align=PP_ALIGN.LEFT, spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Trebuchet MS"
    return tb


def bullets(slide, x, y, w, items, size=16, color=INK, marker="•", marker_color=TEAL):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.4) * len(items))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(11)
        p.line_spacing = 1.2
        r1 = p.add_run()
        r1.text = f"{marker}  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = marker_color
        r1.font.bold = True
        r1.font.name = "Trebuchet MS"
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = "Trebuchet MS"
    return tb


def band(slide, color=TEAL, height=Inches(0.16)):
    bar = slide.shapes.add_shape(1, 0, 0, W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False


def card(slide, x, y, w, h, fill=WHITE, border=LINE):
    sh = slide.shapes.add_shape(5, x, y, w, h)  # rounded rectangle
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.06
    return sh


def heading(slide, title, kicker=None):
    band(slide)
    if kicker:
        textbox(slide, Inches(0.85), Inches(0.55), Inches(11), Inches(0.4),
                kicker.upper(), size=12, color=TEAL, bold=True)
    textbox(slide, Inches(0.85), Inches(0.95), Inches(11.6), Inches(0.9),
            title, size=32, color=INK, bold=True)


def slide_title(prs):
    s = blank(prs)
    hero = s.shapes.add_shape(5, Inches(0.85), Inches(0.9), Inches(11.6), Inches(5.7))
    hero.fill.solid()
    hero.fill.fore_color.rgb = TEAL
    hero.line.fill.background()
    hero.shadow.inherit = False
    hero.adjustments[0] = 0.05

    textbox(s, Inches(1.5), Inches(1.7), Inches(10), Inches(0.5),
            "GRADE 7  ·  PARENT ORIENTATION", size=15, color=WHITE, bold=True)
    textbox(s, Inches(1.5), Inches(2.4), Inches(10.4), Inches(2.2),
            "Jaipur · Abhaneri ·\nRanthambore\nEducational Trip",
            size=44, color=WHITE, bold=True, spacing=1.05)
    textbox(s, Inches(1.5), Inches(5.1), Inches(10), Inches(0.9),
            "Batch 1:  12–19 December 2026        Batch 2:  13–20 December 2026",
            size=17, color=WHITE)
    textbox(s, Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.4),
            "Sample deck — dummy content for demonstration only.",
            size=11, color=MUTED)
    return s


def slide_where(prs):
    s = blank(prs)
    heading(s, "Where the students are going", "Destinations")
    places = [
        ("Jaipur", "Amer Fort, Jantar Mantar and the\nold city's water systems.", TEAL),
        ("Abhaneri", "Chand Baori stepwell — a hands-on\nmeasurement and geometry exercise.", AMBER),
        ("Ranthambore", "Guided safari, habitat mapping and\nfield notes at the camp.", GREEN),
    ]
    x = Inches(0.85)
    for name, desc, colour in places:
        c = card(s, x, Inches(2.3), Inches(3.6), Inches(2.9))
        chip = s.shapes.add_shape(5, x + Inches(0.35), Inches(2.65), Inches(0.55), Inches(0.16))
        chip.fill.solid()
        chip.fill.fore_color.rgb = colour
        chip.line.fill.background()
        chip.shadow.inherit = False
        textbox(s, x + Inches(0.35), Inches(3.0), Inches(3), Inches(0.5),
                name, size=23, color=INK, bold=True)
        textbox(s, x + Inches(0.35), Inches(3.65), Inches(2.95), Inches(1.3),
                desc, size=13.5, color=MUTED)
        x += Inches(3.9)
    return s


def slide_travel(prs):
    s = blank(prs)
    heading(s, "Travel plan", "Getting there and back")
    legs = [
        ("ONWARD", "MMCT JAIPUR SF (12955)", "Departs 10:30 PM",
         "3AC — coach and seat allotment shared a week before.\nReport at the platform 60 minutes before departure.", TEAL),
        ("RETURN", "JP BDTS EXP (12980)", "Departs 8:20 PM",
         "Overnight 3AC.\nArrives Surat the following morning.", AMBER),
    ]
    y = Inches(2.25)
    for tag, train, time, note, colour in legs:
        card(s, Inches(0.85), y, Inches(11.6), Inches(2.0))
        stripe = s.shapes.add_shape(1, Inches(0.85), y, Inches(0.09), Inches(2.0))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = colour
        stripe.line.fill.background()
        stripe.shadow.inherit = False
        textbox(s, Inches(1.3), y + Inches(0.28), Inches(3), Inches(0.35),
                tag, size=12, color=colour, bold=True)
        textbox(s, Inches(1.3), y + Inches(0.66), Inches(6), Inches(0.5),
                train, size=22, color=INK, bold=True)
        textbox(s, Inches(1.3), y + Inches(1.24), Inches(9.5), Inches(0.7),
                note, size=13, color=MUTED)
        textbox(s, Inches(9.0), y + Inches(0.62), Inches(3.2), Inches(0.6),
                time, size=18, color=INK, bold=True, align=PP_ALIGN.RIGHT)
        y += Inches(2.25)
    return s


def slide_itinerary(prs):
    s = blank(prs)
    heading(s, "Day by day", "Outline")
    days = [
        ("Day 1", "Depart Mumbai Central, 10:30 PM"),
        ("Day 2", "Arrive Jaipur · check in · orientation walk"),
        ("Day 3", "Amer Fort, Jantar Mantar, stepwell study"),
        ("Day 4", "Chand Baori measurement exercise, Abhaneri"),
        ("Day 5", "Dawn safari and habitat mapping, Ranthambore"),
        ("Day 6", "Project work and student presentations"),
        ("Day 7", "Board the return train, 8:20 PM"),
    ]
    y = Inches(2.2)
    for i, (day, what) in enumerate(days):
        colour = TEAL if i % 2 == 0 else AMBER
        dot = s.shapes.add_shape(9, Inches(0.9), y + Inches(0.09), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colour
        dot.line.fill.background()
        dot.shadow.inherit = False
        textbox(s, Inches(1.35), y, Inches(1.5), Inches(0.4), day, size=15, color=INK, bold=True)
        textbox(s, Inches(2.9), y, Inches(9.3), Inches(0.4), what, size=15, color=MUTED)
        y += Inches(0.62)
    textbox(s, Inches(0.9), Inches(6.75), Inches(11), Inches(0.4),
            "The full hour-by-hour itinerary is linked on the trip page.", size=12, color=MUTED)
    return s


def slide_safety(prs):
    s = blank(prs)
    heading(s, "Safety and supervision", "How we look after them")
    left = [
        "One accompanying adult for every 12 students.",
        "An adult is with the students at all times, including washroom breaks.",
        "Two dedicated Safety Monitors travel with each batch.",
        "Headcount at every stop and at regular intervals.",
    ]
    right = [
        "A trained medical person and first-aid facility at the campsite at all times.",
        "Nearest police station and hospital details held by the Safety Monitor.",
        "All staff and vendors trained under POCSO, with signed undertakings.",
        "Daily health check at the end of each day, plus WhatsApp updates to parents.",
    ]
    card(s, Inches(0.85), Inches(2.2), Inches(5.65), Inches(4.2))
    card(s, Inches(6.8), Inches(2.2), Inches(5.65), Inches(4.2))
    bullets(s, Inches(1.2), Inches(2.6), Inches(5.0), left, size=14, marker_color=TEAL)
    bullets(s, Inches(7.15), Inches(2.6), Inches(5.0), right, size=14, marker_color=TEAL)
    return s


def slide_dos(prs):
    s = blank(prs)
    heading(s, "Do's and don'ts", "Before you pack")
    card(s, Inches(0.85), Inches(2.2), Inches(5.65), Inches(3.9))
    card(s, Inches(6.8), Inches(2.2), Inches(5.65), Inches(3.9))
    textbox(s, Inches(1.2), Inches(2.5), Inches(4), Inches(0.4), "DO", size=14, color=GREEN, bold=True)
    textbox(s, Inches(7.15), Inches(2.5), Inches(4), Inches(0.4), "DON'T", size=14, color=RED, bold=True)
    bullets(s, Inches(1.2), Inches(3.05), Inches(5.0), [
        "Hand any regular medicine to the Safety Monitor a day before, at the bus stop, with clear written instructions.",
        "Tell the accompanying teacher about any medical history that needs personal attention.",
    ], size=14, marker_color=GREEN)
    bullets(s, Inches(7.15), Inches(3.05), Inches(5.0), [
        "Students may not carry personal medicines — staff carry a first-aid box.",
        "No mobile phones, smart watches or other gadgets. These will be confiscated if found.",
    ], size=14, marker="✕", marker_color=RED)
    return s


def slide_pack(prs):
    s = blank(prs)
    heading(s, "What to pack", "Checklist")
    cols = [
        ["Original School ID and Aadhar card (mandatory)",
         "Clothes for 7 days — full pants, full-sleeve tops",
         "Sweater / fleece jacket / thermals",
         "Heavy woolens, shawl or light blanket",
         "Night dress"],
        ["Sport or trekking shoes (compulsory)",
         "Sandals or slippers",
         "Water bottle and a shoulder bag for it",
         "Towel, napkin, cap, bandana, cotton socks, gloves",
         "Toiletries, moisturizer, sun cream, lip balm"],
    ]
    card(s, Inches(0.85), Inches(2.2), Inches(11.6), Inches(3.6))
    bullets(s, Inches(1.25), Inches(2.6), Inches(5.2), cols[0], size=14, marker="□", marker_color=GREEN)
    bullets(s, Inches(6.9), Inches(2.6), Inches(5.2), cols[1], size=14, marker="□", marker_color=GREEN)
    warn = card(s, Inches(0.85), Inches(6.0), Inches(11.6), Inches(0.85), fill=RGBColor(0xFF, 0xF8, 0xE8),
                border=RGBColor(0xF0, 0xDC, 0xAE))
    textbox(s, Inches(1.25), Inches(6.2), Inches(11), Inches(0.5),
            "Compulsory:  mosquito repellent cream (Odomos)  ·  dry, healthy snacks for the train journey",
            size=14, color=RGBColor(0xB5, 0x79, 0x0C), bold=True)
    return s


def slide_dates(prs):
    s = blank(prs)
    heading(s, "Key dates", "What we need from you")
    items = [("1 September", "Submit signed consent form"),
             ("5 September", "Submit purchase form and trip fee"),
             ("10 September", "Submit medical declaration and ID copies")]
    y = Inches(2.35)
    for date, what in items:
        c = card(s, Inches(0.85), y, Inches(11.6), Inches(1.15), fill=RGBColor(0xFF, 0xF8, 0xE8),
                 border=RGBColor(0xF0, 0xDC, 0xAE))
        stripe = s.shapes.add_shape(1, Inches(0.85), y, Inches(0.09), Inches(1.15))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = AMBER
        stripe.line.fill.background()
        stripe.shadow.inherit = False
        textbox(s, Inches(1.3), y + Inches(0.2), Inches(3), Inches(0.35),
                date.upper(), size=12, color=RGBColor(0xB5, 0x79, 0x0C), bold=True)
        textbox(s, Inches(1.3), y + Inches(0.58), Inches(9.5), Inches(0.45),
                what, size=17, color=INK, bold=True)
        y += Inches(1.4)
    return s


def slide_contact(prs):
    s = blank(prs)
    heading(s, "Who to contact", "Communication")
    rows = [("Trip coordinator", "Ms. Anjali Desai"),
            ("Phone", "+91 98200 11223"),
            ("Email", "trips.grade7@example.edu"),
            ("Emergency (24×7 trip desk)", "+91 98200 99887")]
    card(s, Inches(0.85), Inches(2.25), Inches(11.6), Inches(3.3))
    y = Inches(2.65)
    for label, value in rows:
        textbox(s, Inches(1.3), y, Inches(4.5), Inches(0.4), label.upper(), size=11.5, color=MUTED, bold=True)
        textbox(s, Inches(5.6), y - Inches(0.06), Inches(6.5), Inches(0.45), value, size=17, color=INK, bold=True)
        y += Inches(0.75)
    textbox(s, Inches(0.85), Inches(6.0), Inches(11.6), Inches(0.9),
            "Everything in this deck also lives on the Trip Explorer site, where it stays up to date.\n"
            "Sign in there with the email address or mobile number the school has on record for you.",
            size=14, color=MUTED)
    return s


if __name__ == "__main__":
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = new_deck()
    for fn in (slide_title, slide_where, slide_travel, slide_itinerary,
               slide_safety, slide_dos, slide_pack, slide_dates, slide_contact):
        fn(prs)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

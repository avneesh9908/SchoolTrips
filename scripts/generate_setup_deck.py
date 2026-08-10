"""
Builds the setup deck for the school-management discussion:
what to create, what to name it, what to share, what to send us.

Run:  python scripts/generate_setup_deck.py
"""

import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = pathlib.Path(__file__).resolve().parent.parent
OUT = ROOT / "sample-data" / "Trip Explorer - Setup guide.pptx"

BG = RGBColor(0xFC, 0xF8, 0xED)
INK = RGBColor(0x22, 0x30, 0x3F)
MUTED = RGBColor(0x76, 0x70, 0x66)
TEAL = RGBColor(0x26, 0xC0, 0xB0)
BLUE = RGBColor(0x2A, 0xA8, 0xDE)
BLUE_BG = RGBColor(0xE6, 0xF1, 0xFB)
BLUE_INK = RGBColor(0x0C, 0x44, 0x7C)
AMBER = RGBColor(0xFF, 0xB1, 0x00)
AMBER_BG = RGBColor(0xFF, 0xF8, 0xE8)
AMBER_INK = RGBColor(0xB5, 0x79, 0x0C)
RED = RGBColor(0xFF, 0x6B, 0x5B)
RED_BG = RGBColor(0xFC, 0xEB, 0xEB)
RED_INK = RGBColor(0xA3, 0x2D, 0x2D)
GREEN = RGBColor(0x4C, 0xAF, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xEC, 0xE6, 0xD6)

W, H = Inches(13.333), Inches(7.5)
FONT = "Trebuchet MS"


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(slide, x, y, w, h, body, size=16, color=INK, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.25, mono=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Consolas" if mono else FONT
    return tb


def card(slide, x, y, w, h, fill=WHITE, border=LINE):
    sh = slide.shapes.add_shape(5, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.06
    return sh


def heading(slide, title, kicker=None, accent=TEAL):
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.16))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    if kicker:
        text(slide, Inches(0.85), Inches(0.5), Inches(11), Inches(0.4),
             kicker.upper(), size=12, color=accent, bold=True)
    text(slide, Inches(0.85), Inches(0.9), Inches(11.6), Inches(0.9),
         title, size=31, color=INK, bold=True)


def bullets(slide, x, y, w, items, size=15, marker="•", marker_color=TEAL, color=INK):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.42) * len(items))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.line_spacing = 1.2
        r1 = p.add_run()
        r1.text = f"{marker}  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = marker_color
        r1.font.bold = True
        r1.font.name = FONT
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = FONT
    return tb


def s_title(prs):
    s = blank(prs)
    hero = s.shapes.add_shape(5, Inches(0.85), Inches(0.9), Inches(11.6), Inches(5.7))
    hero.fill.solid()
    hero.fill.fore_color.rgb = TEAL
    hero.line.fill.background()
    hero.shadow.inherit = False
    hero.adjustments[0] = 0.05
    text(s, Inches(1.5), Inches(1.8), Inches(10), Inches(0.5),
         "SETUP GUIDE", size=15, color=WHITE, bold=True)
    text(s, Inches(1.5), Inches(2.5), Inches(10.4), Inches(2.0),
         "Trip Explorer\nWhat we need from you", size=42, color=WHITE, bold=True, spacing=1.1)
    text(s, Inches(1.5), Inches(5.2), Inches(10), Inches(0.8),
         "One Drive folder. Eight spreadsheets. One link.", size=18, color=WHITE)
    return s


def s_overview(prs):
    s = blank(prs)
    heading(s, "The whole setup, in four steps", "Overview")
    steps = [
        ("1", "Name the eight spreadsheets exactly", BLUE_BG, BLUE_INK, BLUE),
        ("2", "Put them in one Drive folder", BLUE_BG, BLUE_INK, BLUE),
        ("3", "Share the folder, copy its link", BLUE_BG, BLUE_INK, BLUE),
        ("4", "Send us that one link", AMBER_BG, AMBER_INK, AMBER),
    ]
    y = Inches(2.2)
    for num, label, fill, ink, accent in steps:
        card(s, Inches(0.85), y, Inches(11.6), Inches(0.95), fill=fill, border=accent)
        chip = s.shapes.add_shape(9, Inches(1.2), y + Inches(0.24), Inches(0.48), Inches(0.48))
        chip.fill.solid()
        chip.fill.fore_color.rgb = accent
        chip.line.fill.background()
        chip.shadow.inherit = False
        text(s, Inches(1.2), y + Inches(0.3), Inches(0.48), Inches(0.4),
             num, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(s, Inches(2.0), y + Inches(0.27), Inches(9.5), Inches(0.5),
             label, size=18, color=ink, bold=True)
        y += Inches(1.12)

    text(s, Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.5),
         "After that, everything you change in the sheets appears on the website. No developer needed.",
         size=14, color=MUTED)
    return s


def s_names(prs):
    s = blank(prs)
    heading(s, "Name the spreadsheets exactly like this", "Step 1", accent=BLUE)
    names = ["Students", "Trips", "Itinerary", "Documents",
             "Guidelines", "Reminders", "Travel", "Media"]
    x, y = Inches(0.85), Inches(2.15)
    for i, n in enumerate(names):
        col, row = i % 4, i // 4
        cx = x + col * Inches(2.95)
        cy = y + row * Inches(1.15)
        card(s, cx, cy, Inches(2.7), Inches(0.95), fill=BLUE_BG, border=BLUE)
        text(s, cx, cy + Inches(0.28), Inches(2.7), Inches(0.5),
             n, size=17, color=BLUE_INK, bold=True, align=PP_ALIGN.CENTER)

    card(s, Inches(0.85), Inches(4.7), Inches(11.6), Inches(1.9), fill=RED_BG, border=RED)
    text(s, Inches(1.25), Inches(4.95), Inches(11), Inches(0.4),
         "The name is how we find each sheet", size=15, color=RED_INK, bold=True)
    bullets(s, Inches(1.25), Inches(5.4), Inches(10.8), [
        "All plural. Never Student, Document or Reminder.",
        "Extra words are fine — \"Grade 7 Itinerary 2026\" still works. Two files matching the same name do not.",
        "Each must be a Google Sheet. An uploaded .xlsx must be opened and saved as a Google Sheet.",
    ], size=13.5, marker_color=RED, color=INK)
    return s


def s_folder(prs):
    s = blank(prs)
    heading(s, "Put everything in one folder", "Step 2", accent=BLUE)
    card(s, Inches(0.85), Inches(2.15), Inches(5.6), Inches(4.3))
    text(s, Inches(1.25), Inches(2.45), Inches(5), Inches(0.4),
         "School Trips", size=17, color=INK, bold=True)
    text(s, Inches(1.5), Inches(3.0), Inches(5), Inches(3.2),
         "Students\nTrips\nItinerary\nDocuments\nGuidelines\nReminders\nTravel\nMedia",
         size=14, color=MUTED, mono=True, spacing=1.35)

    card(s, Inches(6.85), Inches(2.15), Inches(5.6), Inches(4.3))
    text(s, Inches(7.25), Inches(2.45), Inches(5), Inches(0.4),
         "…and a folder per grade", size=17, color=INK, bold=True)
    text(s, Inches(7.5), Inches(3.0), Inches(5), Inches(2.4),
         "Grade 7\nGrade 9\nGrade 5", size=14, color=MUTED, mono=True, spacing=1.35)
    text(s, Inches(7.25), Inches(4.6), Inches(5), Inches(1.6),
         "Orientation decks, posters and photos go\nin these. Link them from the Documents\nsheet — or paste the folder link once and\nevery file in it appears.",
         size=13, color=MUTED)
    return s


def s_sharing(prs):
    s = blank(prs)
    heading(s, "Sharing — please treat as two jobs", "Step 3", accent=AMBER)
    items = [
        ("The folder and every spreadsheet in it",
         "Anyone with the link → Viewer. Without this the website reads nothing at all."),
        ("Each document you want previewed",
         "Same setting. A private document still opens for those who have access — it just shows a plain tile instead of a picture."),
    ]
    y = Inches(2.2)
    for title, body in items:
        card(s, Inches(0.85), y, Inches(11.6), Inches(1.55), fill=AMBER_BG, border=AMBER)
        text(s, Inches(1.3), y + Inches(0.25), Inches(10.8), Inches(0.4),
             title, size=17, color=AMBER_INK, bold=True)
        text(s, Inches(1.3), y + Inches(0.75), Inches(10.6), Inches(0.7),
             body, size=13.5, color=INK)
        y += Inches(1.85)

    card(s, Inches(0.85), Inches(5.95), Inches(11.6), Inches(1.0), fill=RED_BG, border=RED)
    text(s, Inches(1.3), Inches(6.2), Inches(10.8), Inches(0.6),
         "Before using real family data, please read the privacy note we sent — the student list becomes publicly readable.",
         size=14, color=RED_INK, bold=True)
    return s


def s_mistakes(prs):
    s = blank(prs)
    heading(s, "What quietly breaks it", "Worth knowing", accent=RED)
    rows = [
        ("A sheet renamed", "That whole section disappears, with no error anywhere."),
        ("A grade spelled differently", "\"Grade Seven\" or \"7th std\" — that row vanishes. Use the dropdown."),
        ("An .xlsx never converted", "Upload alone is not enough. Open it, File → Save as Google Sheets."),
        ("Two files with similar names", "\"Winter Trips\" and \"Summer Trips\" both match Trips, so neither is used."),
        ("A parent with no email and no phone", "That child becomes invisible to everyone."),
    ]
    y = Inches(2.15)
    for what, effect in rows:
        card(s, Inches(0.85), y, Inches(11.6), Inches(0.86))
        stripe = s.shapes.add_shape(1, Inches(0.85), y, Inches(0.08), Inches(0.86))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RED
        stripe.line.fill.background()
        stripe.shadow.inherit = False
        text(s, Inches(1.25), y + Inches(0.24), Inches(3.6), Inches(0.4),
             what, size=14.5, color=INK, bold=True)
        text(s, Inches(5.0), y + Inches(0.26), Inches(7.2), Inches(0.4),
             effect, size=13, color=MUTED)
        y += Inches(0.98)
    return s


def s_send(prs):
    s = blank(prs)
    heading(s, "What to send us", "Step 4", accent=TEAL)
    card(s, Inches(0.85), Inches(2.15), Inches(11.6), Inches(2.5), fill=RGBColor(0xE1, 0xF5, 0xEE), border=TEAL)
    text(s, Inches(1.3), Inches(2.5), Inches(10.8), Inches(0.5),
         "One link", size=22, color=RGBColor(0x0F, 0x6E, 0x56), bold=True)
    text(s, Inches(1.3), Inches(3.15), Inches(10.8), Inches(0.6),
         "https://drive.google.com/drive/folders/…", size=15, color=RGBColor(0x1D, 0x9E, 0x75), mono=True)
    text(s, Inches(1.3), Inches(3.8), Inches(10.8), Inches(0.6),
         "That is the entire configuration. Everything else is reached through it.",
         size=14, color=INK)

    bullets(s, Inches(1.25), Inches(5.0), Inches(11), [
        "Confirmation that the folder and its sheets are shared with Anyone with the link → Viewer.",
        "Which grades are in scope for the first launch.",
        "Your answers to the open questions: privacy, how parents sign in, and whether \"Father\" should read \"Parent\".",
    ], size=14, marker_color=TEAL)
    return s


if __name__ == "__main__":
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for fn in (s_title, s_overview, s_names, s_folder, s_sharing, s_mistakes, s_send):
        fn(prs)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)} ({len(prs.slides._sldIdLst)} slides)")
